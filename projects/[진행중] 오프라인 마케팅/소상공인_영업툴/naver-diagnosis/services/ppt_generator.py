"""
PPT 제안서 생성기
python-pptx를 사용하여 9페이지 진단 보고서 생성 (경쟁사 비교 + 손익분기 슬라이드 추가)
"""
import math
import os
from datetime import datetime
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class PPTGenerator:
    """PPT 제안서 생성기"""

    # 색상 정의
    COLORS = {
        "primary": RGBColor(0x00, 0x5A, 0x9E),      # 네이버 블루
        "secondary": RGBColor(0x03, 0xC7, 0x5A),    # 네이버 그린
        "dark": RGBColor(0x33, 0x33, 0x33),
        "gray": RGBColor(0x66, 0x66, 0x66),
        "light_gray": RGBColor(0xF5, 0xF5, 0xF5),
        "grade_a": RGBColor(0x00, 0xC7, 0x3C),      # 녹색
        "grade_b": RGBColor(0x4C, 0xAF, 0x50),      # 연녹색
        "grade_c": RGBColor(0xFF, 0x98, 0x00),      # 주황색
        "grade_d": RGBColor(0xF4, 0x43, 0x36),      # 빨간색
    }

    # 한글 폰트 (Windows 기본)
    FONT_NAME = "맑은 고딕"

    def __init__(self, output_dir: str = "ppt_output"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def _get_grade_color(self, grade: str) -> RGBColor:
        """등급별 색상 반환"""
        colors = {
            "A": self.COLORS["grade_a"],
            "B": self.COLORS["grade_b"],
            "C": self.COLORS["grade_c"],
            "D": self.COLORS["grade_d"],
        }
        return colors.get(grade, self.COLORS["gray"])

    def _add_text_frame(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int = 18,
        bold: bool = False,
        color: Optional[RGBColor] = None,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        """텍스트 프레임 추가"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONT_NAME
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment
        if color:
            p.font.color.rgb = color
        return shape

    def _create_cover_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 1: 표지 — 충격 문구 포함"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # 배경 색상 박스 (상단)
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS["primary"]
        shape.line.fill.background()

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.8, 9, 1,
            "네이버 플레이스 진단 보고서",
            font_size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        # 업체명
        self._add_text_frame(
            slide, 0.5, 2.8, 9, 0.8,
            data.get("business_name", "업체명"),
            font_size=32, bold=True, color=self.COLORS["dark"],
            alignment=PP_ALIGN.CENTER
        )

        # 분석일
        analysis_date = datetime.now().strftime("%Y년 %m월 %d일")
        self._add_text_frame(
            slide, 0.5, 3.6, 9, 0.4,
            f"분석일: {analysis_date}",
            font_size=14, color=self.COLORS["gray"],
            alignment=PP_ALIGN.CENTER
        )

        # 충격 문구 박스 — 손실 고객 수 (있을 때만)
        estimated_lost = data.get("estimated_lost_customers", 0)
        if estimated_lost and estimated_lost > 0:
            shock_box = slide.shapes.add_shape(
                1, Inches(1.0), Inches(4.1), Inches(8), Inches(0.7)
            )
            shock_box.fill.solid()
            shock_box.fill.fore_color.rgb = self.COLORS["grade_d"]
            shock_box.line.fill.background()

            self._add_text_frame(
                slide, 1.0, 4.15, 8, 0.6,
                f"사장님, 지금 매달 약 {estimated_lost}명이 경쟁사로 가고 있어요",
                font_size=17, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER
            )

        # 종합점수 박스
        score_y = 5.0 if estimated_lost and estimated_lost > 0 else 4.5
        score_shape = slide.shapes.add_shape(
            5, Inches(3.5), Inches(score_y), Inches(3), Inches(1.5)
        )
        score_shape.fill.solid()
        score_shape.fill.fore_color.rgb = self._get_grade_color(data.get("grade", "D"))
        score_shape.line.fill.background()

        # 종합점수 텍스트
        self._add_text_frame(
            slide, 3.5, score_y + 0.1, 3, 0.5,
            "종합점수",
            font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        # 점수 및 등급
        total_score = data.get("total_score", 0)
        grade = data.get("grade", "D")
        self._add_text_frame(
            slide, 3.5, score_y + 0.6, 3, 0.8,
            f"{total_score:.0f}점 (등급 {grade})",
            font_size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

    def _create_analysis_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 2: 현황 분석 (7개 항목)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "현황 분석",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        # 항목별 점수 표시 (7개로 확장)
        items = [
            ("사진", "photo_score", data.get("photo_count", 0), "장"),
            ("리뷰 (영수증+방문자)", "review_score", data.get("review_count", 0), "개"),
            ("블로그 리뷰", "blog_score", data.get("blog_review_count", 0), "개"),
            ("정보 완성도", "info_score", None, None),
            ("키워드 노출", "keyword_score", len(data.get("keywords", [])), "개"),
            ("편의 기능", "convenience_score", None, None),
            ("고객 참여", "engagement_score", None, None),
        ]

        y_pos = 1.3
        for label, score_key, count, unit in items:
            score = data.get(score_key, 0)

            # 항목명
            self._add_text_frame(
                slide, 0.5, y_pos, 3, 0.5,
                label,
                font_size=14, bold=True, color=self.COLORS["dark"]
            )

            # 점수 바 (배경)
            bar_bg = slide.shapes.add_shape(
                1, Inches(3.5), Inches(y_pos + 0.05), Inches(5), Inches(0.3)
            )
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = self.COLORS["light_gray"]
            bar_bg.line.fill.background()

            # 점수 바 (실제)
            bar_width = max(0.1, (score / 100) * 5)
            bar = slide.shapes.add_shape(
                1, Inches(3.5), Inches(y_pos + 0.05), Inches(bar_width), Inches(0.3)
            )
            bar.fill.solid()
            if score >= 80:
                bar.fill.fore_color.rgb = self.COLORS["grade_a"]
            elif score >= 60:
                bar.fill.fore_color.rgb = self.COLORS["grade_b"]
            elif score >= 40:
                bar.fill.fore_color.rgb = self.COLORS["grade_c"]
            else:
                bar.fill.fore_color.rgb = self.COLORS["grade_d"]
            bar.line.fill.background()

            # 점수 텍스트
            score_text = f"{score:.0f}점"
            if count is not None and unit:
                score_text = f"{score:.0f}점 ({count}{unit})"

            self._add_text_frame(
                slide, 8.7, y_pos, 1, 0.5,
                score_text,
                font_size=12, color=self.COLORS["gray"],
                alignment=PP_ALIGN.RIGHT
            )

            y_pos += 0.75  # 간격을 0.75로 줄여서 7개가 들어가도록

        # 정보 완성도 상세
        y_pos += 0.15
        self._add_text_frame(
            slide, 0.5, y_pos, 9, 0.4,
            "정보 등록 현황",
            font_size=12, bold=True, color=self.COLORS["gray"]
        )

        y_pos += 0.4
        info_items = [
            ("영업시간", data.get("has_hours", False)),
            ("메뉴", data.get("has_menu", False)),
            ("가격정보", data.get("has_price", False)),
            ("소개글", data.get("has_intro", False)),
            ("오시는 길", data.get("has_directions", False)),
        ]

        # 2열로 배치
        col1_items = info_items[:3]
        col2_items = info_items[3:]

        temp_y = y_pos
        for item_name, has_item in col1_items:
            status = "등록" if has_item else "미등록"
            status_color = self.COLORS["grade_a"] if has_item else self.COLORS["grade_d"]

            self._add_text_frame(
                slide, 0.7, temp_y, 1.5, 0.35,
                f"- {item_name}:",
                font_size=11, color=self.COLORS["dark"]
            )
            self._add_text_frame(
                slide, 2.2, temp_y, 1, 0.35,
                status,
                font_size=11, bold=True, color=status_color
            )
            temp_y += 0.3

        temp_y = y_pos
        for item_name, has_item in col2_items:
            status = "등록" if has_item else "미등록"
            status_color = self.COLORS["grade_a"] if has_item else self.COLORS["grade_d"]

            self._add_text_frame(
                slide, 5.0, temp_y, 1.5, 0.35,
                f"- {item_name}:",
                font_size=11, color=self.COLORS["dark"]
            )
            self._add_text_frame(
                slide, 6.5, temp_y, 1, 0.35,
                status,
                font_size=11, bold=True, color=status_color
            )
            temp_y += 0.3

    def _create_detail_check_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 3: 서비스 세팅 체크리스트 (신규)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "서비스 세팅 체크리스트",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        # 플레이스 순위 정보 박스
        rank = data.get("naver_place_rank", 0)
        rank_text = f"현재 플레이스 순위: {rank}위" if rank > 0 else "현재 순위: 확인 중"
        rank_color = self.COLORS["grade_a"] if rank > 0 and rank <= 10 else (self.COLORS["grade_c"] if rank > 0 else self.COLORS["gray"])

        rank_box = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6)
        )
        rank_box.fill.solid()
        rank_box.fill.fore_color.rgb = self.COLORS["light_gray"]
        rank_box.line.fill.background()

        self._add_text_frame(
            slide, 0.5, 1.3, 9, 0.5,
            rank_text,
            font_size=18, bold=True, color=rank_color,
            alignment=PP_ALIGN.CENTER
        )

        # 왼쪽 컬럼 - 편의기능
        self._add_text_frame(
            slide, 0.5, 2.0, 4, 0.5,
            "편의 기능",
            font_size=16, bold=True, color=self.COLORS["dark"]
        )

        left_items = [
            ("네이버 예약", data.get("has_booking", False)),
            ("톡톡 상담", data.get("has_talktalk", False)),
            ("스마트콜", data.get("has_smartcall", False)),
            ("쿠폰/이벤트", data.get("has_coupon", False)),
            ("새소식", data.get("has_news", False)),
        ]

        y_pos = 2.5
        for item_name, has_item in left_items:
            status = "등록" if has_item else "미등록"
            status_color = self.COLORS["grade_a"] if has_item else self.COLORS["grade_d"]

            # 항목 배경 박스
            item_box = slide.shapes.add_shape(
                1, Inches(0.5), Inches(y_pos), Inches(4), Inches(0.45)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            item_box.line.fill.background()

            self._add_text_frame(
                slide, 0.6, y_pos + 0.08, 2.5, 0.35,
                item_name,
                font_size=13, color=self.COLORS["dark"]
            )
            self._add_text_frame(
                slide, 3.3, y_pos + 0.08, 1, 0.35,
                status,
                font_size=13, bold=True, color=status_color,
                alignment=PP_ALIGN.RIGHT
            )
            y_pos += 0.55

        # 오른쪽 컬럼 - 콘텐츠/채널
        self._add_text_frame(
            slide, 5.0, 2.0, 4, 0.5,
            "콘텐츠 / 외부 채널",
            font_size=16, bold=True, color=self.COLORS["dark"]
        )

        right_items = [
            ("메뉴 상세설명", data.get("has_menu_description", False)),
            ("사장님 답글", data.get("has_owner_reply", False)),
            ("인스타그램 연동", data.get("has_instagram", False)),
            ("카카오 채널", data.get("has_kakao", False)),
        ]

        y_pos = 2.5
        for item_name, has_item in right_items:
            status = "등록" if has_item else "미등록"
            status_color = self.COLORS["grade_a"] if has_item else self.COLORS["grade_d"]

            # 항목 배경 박스
            item_box = slide.shapes.add_shape(
                1, Inches(5.0), Inches(y_pos), Inches(4.5), Inches(0.45)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
            item_box.line.fill.background()

            self._add_text_frame(
                slide, 5.1, y_pos + 0.08, 2.8, 0.35,
                item_name,
                font_size=13, color=self.COLORS["dark"]
            )
            self._add_text_frame(
                slide, 8.2, y_pos + 0.08, 1, 0.35,
                status,
                font_size=13, bold=True, color=status_color,
                alignment=PP_ALIGN.RIGHT
            )
            y_pos += 0.55

        # 메뉴 수 표시 (숫자)
        menu_count = data.get("menu_count", 0)
        item_box = slide.shapes.add_shape(
            1, Inches(5.0), Inches(y_pos), Inches(4.5), Inches(0.45)
        )
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        item_box.line.fill.background()

        self._add_text_frame(
            slide, 5.1, y_pos + 0.08, 2.8, 0.35,
            "등록된 메뉴 수",
            font_size=13, color=self.COLORS["dark"]
        )
        menu_color = self.COLORS["grade_a"] if menu_count >= 10 else (self.COLORS["grade_c"] if menu_count > 0 else self.COLORS["grade_d"])
        self._add_text_frame(
            slide, 8.2, y_pos + 0.08, 1, 0.35,
            f"{menu_count}개",
            font_size=13, bold=True, color=menu_color,
            alignment=PP_ALIGN.RIGHT
        )

        # 하단 안내 문구
        self._add_text_frame(
            slide, 0.5, 6.8, 9, 0.5,
            "* 편의 기능과 외부 채널 연동은 고객 유입과 전환율에 큰 영향을 미칩니다.",
            font_size=11, color=self.COLORS["gray"]
        )

    def _create_improvement_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 6: 비어있는 항목 리스트 (개선 방법 표시 안 함)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "지금 비어있는 항목들",
            font_size=28, bold=True, color=self.COLORS["grade_d"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["grade_d"]
        line.line.fill.background()

        improvement_points = data.get("improvement_points", [])

        if not improvement_points:
            self._add_text_frame(
                slide, 0.5, 2.5, 9, 1,
                "현재 플레이스 항목이 모두 등록되어 있습니다.",
                font_size=18, color=self.COLORS["grade_a"],
                alignment=PP_ALIGN.CENTER
            )
            return

        # 안내 문구
        self._add_text_frame(
            slide, 0.5, 1.1, 9, 0.4,
            "아래 항목들이 비어있어 검색 노출에서 손해를 보고 있어요.",
            font_size=14, color=self.COLORS["gray"]
        )

        y_pos = 1.65
        for i, point in enumerate(improvement_points[:5], 1):
            category = point.get("category", "")
            message = point.get("message", "")

            category_labels = {
                "photo": "사진",
                "review": "리뷰",
                "blog": "블로그",
                "info": "기본정보",
                "keyword": "키워드",
                "convenience": "편의기능",
                "engagement": "고객소통",
            }
            category_label = category_labels.get(category, category)

            # X 마크 박스
            x_box = slide.shapes.add_shape(
                1, Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.45)
            )
            x_box.fill.solid()
            x_box.fill.fore_color.rgb = self.COLORS["grade_d"]
            x_box.line.fill.background()

            self._add_text_frame(
                slide, 0.5, y_pos + 0.05, 0.5, 0.35,
                "✕",
                font_size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER
            )

            # 카테고리 뱃지
            self._add_text_frame(
                slide, 1.15, y_pos + 0.05, 1.2, 0.35,
                f"[{category_label}]",
                font_size=12, bold=True, color=self.COLORS["grade_d"]
            )

            # 메시지 (비어있는 항목만, 해결책 없음)
            self._add_text_frame(
                slide, 2.5, y_pos + 0.05, 7, 0.4,
                message,
                font_size=13, color=self.COLORS["dark"]
            )

            y_pos += 0.7

        # 하단 CTA
        self._add_text_frame(
            slide, 0.5, 6.8, 9, 0.4,
            "이 항목들, 저희가 전부 채워드릴 수 있어요.",
            font_size=13, bold=True, color=self.COLORS["primary"],
            alignment=PP_ALIGN.CENTER
        )

    def _create_keyword_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 4: 키워드 & 순위 현황"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "키워드 검색량 & 순위 현황",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        # 순위 정보 박스 (상단)
        rank = data.get("naver_place_rank", 0)
        rank_text = f"현재 플레이스 순위: {rank}위" if rank > 0 else "현재 순위: 확인 중"
        rank_color = self.COLORS["grade_a"] if rank > 0 and rank <= 10 else (self.COLORS["grade_c"] if rank > 0 else self.COLORS["gray"])

        rank_box = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.15), Inches(9), Inches(0.5)
        )
        rank_box.fill.solid()
        rank_box.fill.fore_color.rgb = self.COLORS["light_gray"]
        rank_box.line.fill.background()

        self._add_text_frame(
            slide, 0.5, 1.22, 9, 0.4,
            rank_text,
            font_size=16, bold=True, color=rank_color,
            alignment=PP_ALIGN.CENTER
        )

        keywords = data.get("keywords", [])

        if not keywords:
            self._add_text_frame(
                slide, 0.5, 3.0, 9, 1,
                "키워드 데이터 수집 예정\n\n타겟 키워드를 설정하면 월간 검색량 분석이 가능합니다.",
                font_size=16, color=self.COLORS["gray"],
                alignment=PP_ALIGN.CENTER
            )
            return

        # 테이블 헤더
        y_pos = 1.8
        headers = [("키워드", 0.5, 3.5), ("월간 검색량", 4.2, 1.8), ("순위", 6.2, 1.3), ("경쟁도", 7.7, 1.8)]
        for header, x, width in headers:
            header_shape = slide.shapes.add_shape(
                1, Inches(x), Inches(y_pos), Inches(width), Inches(0.45)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = self.COLORS["primary"]
            header_shape.line.fill.background()

            self._add_text_frame(
                slide, x, y_pos + 0.08, width, 0.35,
                header,
                font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER
            )

        y_pos += 0.55

        # 데이터 행
        for kw in keywords[:8]:  # 최대 8개
            # 문자열 또는 딕셔너리 형태 모두 처리
            if isinstance(kw, str):
                keyword = kw
                volume = 0
                kw_rank = 0
            else:
                keyword = kw.get("keyword", "")
                volume = kw.get("search_volume", 0)
                kw_rank = kw.get("rank", 0)

            # 키워드
            self._add_text_frame(
                slide, 0.6, y_pos, 3.3, 0.4,
                keyword,
                font_size=12, color=self.COLORS["dark"]
            )

            # 검색량
            volume_text = f"{volume:,}" if volume > 0 else "-"
            self._add_text_frame(
                slide, 4.2, y_pos, 1.8, 0.4,
                volume_text,
                font_size=12, color=self.COLORS["dark"],
                alignment=PP_ALIGN.CENTER
            )

            # 순위
            rank_display = f"{kw_rank}위" if kw_rank > 0 else "-"
            rank_display_color = self.COLORS["grade_a"] if kw_rank > 0 and kw_rank <= 10 else (self.COLORS["grade_c"] if kw_rank > 0 else self.COLORS["gray"])
            self._add_text_frame(
                slide, 6.2, y_pos, 1.3, 0.4,
                rank_display,
                font_size=12, bold=True, color=rank_display_color,
                alignment=PP_ALIGN.CENTER
            )

            # 경쟁도 (검색량 기반 추정)
            if volume >= 10000:
                competition = "높음"
                comp_color = self.COLORS["grade_d"]
            elif volume >= 1000:
                competition = "보통"
                comp_color = self.COLORS["grade_c"]
            else:
                competition = "낮음"
                comp_color = self.COLORS["grade_a"]

            self._add_text_frame(
                slide, 7.7, y_pos, 1.8, 0.4,
                competition,
                font_size=12, bold=True, color=comp_color,
                alignment=PP_ALIGN.CENTER
            )

            # 구분선
            sep_line = slide.shapes.add_shape(
                1, Inches(0.5), Inches(y_pos + 0.4), Inches(9), Inches(0.01)
            )
            sep_line.fill.solid()
            sep_line.fill.fore_color.rgb = self.COLORS["light_gray"]
            sep_line.line.fill.background()

            y_pos += 0.5

    def _create_related_keywords_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 5: 연관 키워드 제안 (신규)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "연관 키워드 제안",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        related_keywords = data.get("related_keywords", [])

        if not related_keywords:
            self._add_text_frame(
                slide, 0.5, 2.5, 9, 1,
                "연관 키워드 데이터가 없습니다.\n\n업종 및 지역 기반 연관 키워드를 분석 중입니다.",
                font_size=16, color=self.COLORS["gray"],
                alignment=PP_ALIGN.CENTER
            )
            return

        # 설명 텍스트
        self._add_text_frame(
            slide, 0.5, 1.15, 9, 0.5,
            "아래 키워드를 메뉴명, 소개글, 해시태그에 활용하면 검색 노출이 향상됩니다.",
            font_size=14, color=self.COLORS["gray"]
        )

        # 키워드들을 2열로 배치
        y_pos = 1.7
        for i, kw in enumerate(related_keywords[:10]):
            x_pos = 0.5 if i % 2 == 0 else 5.0
            if i % 2 == 0 and i > 0:
                y_pos += 0.6

            # 키워드 박스
            box = slide.shapes.add_shape(
                1, Inches(x_pos), Inches(y_pos), Inches(4.3), Inches(0.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS["light_gray"]
            box.line.fill.background()

            # 키워드가 딕셔너리인 경우와 문자열인 경우 처리
            if isinstance(kw, dict):
                kw_text = kw.get("keyword", str(kw))
            else:
                kw_text = str(kw)

            self._add_text_frame(
                slide, x_pos + 0.15, y_pos + 0.08, 4, 0.35,
                f"# {kw_text}",
                font_size=13, color=self.COLORS["primary"]
            )

        # 안내 텍스트
        self._add_text_frame(
            slide, 0.5, 6.8, 9, 0.5,
            "* 연관 키워드를 활용한 콘텐츠 작성 시 자연스러운 문맥으로 포함시키는 것이 효과적입니다.",
            font_size=11, color=self.COLORS["gray"]
        )

    def _create_competitor_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 (신규): 경쟁사 비교 — 내 업체 vs 지역 평균"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "경쟁사 비교",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        # 안내 문구
        rank = data.get("naver_place_rank", 0)
        if rank > 0:
            self._add_text_frame(
                slide, 0.5, 1.1, 9, 0.4,
                f"현재 검색 순위 {rank}위 — 1~5위 밖이면 사실상 안 보여요",
                font_size=14, color=self.COLORS["grade_d"] if rank > 5 else self.COLORS["gray"]
            )

        # 데이터
        my_review = data.get("review_count", 0)
        my_photo = data.get("photo_count", 0)
        my_blog = data.get("blog_review_count", 0)
        competitor_avg_review = data.get("competitor_avg_review", 0)
        competitor_avg_photo = data.get("competitor_avg_photo", 0)
        competitor_avg_blog = data.get("competitor_avg_blog", 0)

        # 비교 항목 목록
        compare_items = [
            ("리뷰 수", my_review, competitor_avg_review, "개"),
            ("사진 수", my_photo, competitor_avg_photo, "장"),
            ("블로그 리뷰", my_blog, competitor_avg_blog, "개"),
        ]

        # 헤더 배경
        header_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.65), Inches(9), Inches(0.45))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.COLORS["primary"]
        header_box.line.fill.background()

        # 헤더 텍스트
        self._add_text_frame(
            slide, 0.6, 1.7, 3, 0.35,
            "항목",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
        )
        self._add_text_frame(
            slide, 3.8, 1.7, 2.5, 0.35,
            "우리 업체",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )
        self._add_text_frame(
            slide, 6.4, 1.7, 2.8, 0.35,
            "지역 상위 평균",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        y_pos = 2.2
        for label, my_val, comp_val, unit in compare_items:
            # 행 배경 (교차 색상)
            row_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
            row_box.fill.solid()
            row_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
            row_box.line.fill.background()

            # 항목명
            self._add_text_frame(
                slide, 0.6, y_pos + 0.08, 3, 0.35,
                label,
                font_size=13, color=self.COLORS["dark"]
            )

            # 내 값
            my_color = self.COLORS["grade_a"] if comp_val <= 0 or my_val >= comp_val else self.COLORS["grade_d"]
            self._add_text_frame(
                slide, 3.8, y_pos + 0.08, 2.5, 0.35,
                f"{my_val:,}{unit}" if my_val > 0 else f"0{unit}",
                font_size=14, bold=True, color=my_color,
                alignment=PP_ALIGN.CENTER
            )

            # 경쟁사 값
            self._add_text_frame(
                slide, 6.4, y_pos + 0.08, 2.8, 0.35,
                f"{comp_val:,}{unit}" if comp_val > 0 else "데이터 없음",
                font_size=14, bold=True, color=self.COLORS["dark"],
                alignment=PP_ALIGN.CENTER
            )

            y_pos += 0.6

        # 손실 고객 강조 박스
        estimated_lost = data.get("estimated_lost_customers", 0)
        if estimated_lost and estimated_lost > 0:
            lost_box = slide.shapes.add_shape(
                1, Inches(0.5), Inches(y_pos + 0.3), Inches(9), Inches(0.8)
            )
            lost_box.fill.solid()
            lost_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF3)
            lost_box.line.fill.background()

            self._add_text_frame(
                slide, 0.7, y_pos + 0.4, 8.5, 0.6,
                f"⚠️  현재 순위 기준, 매달 약 {estimated_lost}명이 경쟁사로 가고 있는 걸로 추정됩니다.",
                font_size=15, bold=True, color=self.COLORS["grade_d"]
            )

    def _create_breakeven_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 (신규): 손익분기 슬라이드"""
        from config.industry_weights import get_avg_price, PACKAGES, recommend_package

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목
        self._add_text_frame(
            slide, 0.5, 0.3, 9, 0.7,
            "손익분기 계산",
            font_size=28, bold=True, color=self.COLORS["primary"]
        )

        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["primary"]
        line.line.fill.background()

        category = data.get("category", "")
        avg_price = get_avg_price(category)

        # 업종 평균 객단가 표시
        self._add_text_frame(
            slide, 0.5, 1.15, 9, 0.4,
            f"업종 평균 객단가: {avg_price:,}원 기준",
            font_size=14, color=self.COLORS["gray"]
        )

        # 패키지별 손익분기 계산
        packages_info = [
            ("주목", 290000),
            ("집중", 490000),
            ("시선", 890000),
        ]

        # 추천 패키지 계산
        from services.scorer import DiagnosisScorer
        weak = []
        for key in ["photo", "review", "blog", "info", "keyword", "convenience", "engagement"]:
            if data.get(f"{key}_score", 100) < 50:
                weak.append(key)
        grade = data.get("grade", "D")
        recommended = recommend_package(grade, weak)

        # 헤더
        header_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.65), Inches(9), Inches(0.45))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.COLORS["primary"]
        header_box.line.fill.background()

        self._add_text_frame(
            slide, 0.6, 1.7, 3, 0.35,
            "패키지",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
        )
        self._add_text_frame(
            slide, 3.8, 1.7, 2.5, 0.35,
            "월 비용",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )
        self._add_text_frame(
            slide, 6.4, 1.7, 2.8, 0.35,
            "손익분기 고객 수",
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        pkg_labels = {"주목": "주목 플랜", "집중": "집중 플랜", "시선": "시선 플랜"}
        y_pos = 2.2
        for pkg_name, pkg_price in packages_info:
            breakeven = math.ceil(pkg_price / avg_price) if avg_price > 0 else 0
            is_recommended = (pkg_name == recommended)

            # 행 배경 (추천 패키지는 강조)
            row_color = RGBColor(0xFF, 0xFB, 0xE8) if is_recommended else RGBColor(0xF8, 0xF8, 0xF8)
            row_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.55))
            row_box.fill.solid()
            row_box.fill.fore_color.rgb = row_color
            row_box.line.fill.background()

            label = f"{'★ ' if is_recommended else ''}{pkg_labels[pkg_name]}"
            label_color = self.COLORS["secondary"] if is_recommended else self.COLORS["dark"]

            self._add_text_frame(
                slide, 0.6, y_pos + 0.1, 3, 0.38,
                label,
                font_size=13, bold=is_recommended, color=label_color
            )
            self._add_text_frame(
                slide, 3.8, y_pos + 0.1, 2.5, 0.38,
                f"{pkg_price:,}원/월",
                font_size=13, color=self.COLORS["dark"],
                alignment=PP_ALIGN.CENTER
            )

            breakeven_text = f"{breakeven}명만 더" if breakeven > 0 else "-"
            self._add_text_frame(
                slide, 6.4, y_pos + 0.1, 2.8, 0.38,
                breakeven_text,
                font_size=14, bold=is_recommended, color=self.COLORS["grade_a"] if breakeven <= 5 else self.COLORS["dark"],
                alignment=PP_ALIGN.CENTER
            )

            y_pos += 0.7

        # 강조 메시지
        estimated_lost = data.get("estimated_lost_customers", 0)
        if estimated_lost and estimated_lost > 0:
            msg = f"지금 매달 {estimated_lost}명이 경쟁사로 가고 있어요. 손해 보고 계신 금액이 더 크죠."
        else:
            msg = "고객 몇 명만 더 와도 패키지 비용은 충분히 회수됩니다."

        self._add_text_frame(
            slide, 0.5, y_pos + 0.3, 9, 0.5,
            msg,
            font_size=14, bold=True, color=self.COLORS["primary"],
            alignment=PP_ALIGN.CENTER
        )

    def _create_proposal_slide(self, prs: Presentation, data: Dict[str, Any]):
        """페이지 7: 제안"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 배경 색상 박스 (상단)
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS["primary"]
        shape.line.fill.background()

        # 제목
        self._add_text_frame(
            slide, 0.5, 1.5, 9, 1,
            "플레이스 최적화 제안",
            font_size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        # 메인 메시지
        grade = data.get("grade", "D")
        total_score = data.get("total_score", 0)

        if grade == "A":
            main_message = "현재 플레이스 관리 상태가 우수합니다.\n지속적인 관리로 고객 유입을 유지하세요."
        elif grade == "B":
            main_message = "플레이스 관리 상태가 양호합니다.\n몇 가지 개선으로 더 많은 고객을 유치할 수 있습니다."
        elif grade == "C":
            main_message = "플레이스 최적화가 필요합니다.\n체계적인 관리로 검색 노출을 높일 수 있습니다."
        else:
            main_message = "플레이스 기본 설정이 부족합니다.\n지금 바로 개선을 시작하세요."

        self._add_text_frame(
            slide, 0.5, 2.8, 9, 1.5,
            main_message,
            font_size=20, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        # 핵심 제안
        self._add_text_frame(
            slide, 0.5, 4.5, 9, 0.8,
            "플레이스 최적화를 통해 고객 유입을 늘릴 수 있습니다",
            font_size=18, bold=True, color=self.COLORS["secondary"],
            alignment=PP_ALIGN.CENTER
        )

        # 예상 효과
        expected_improvement = {
            "D": "최대 200%",
            "C": "최대 100%",
            "B": "최대 50%",
            "A": "현재 유지",
        }

        self._add_text_frame(
            slide, 0.5, 5.3, 9, 0.6,
            f"예상 고객 유입 증가: {expected_improvement.get(grade, '측정 필요')}",
            font_size=16, color=RGBColor(0xFF, 0xFF, 0xFF),
            alignment=PP_ALIGN.CENTER
        )

        # 연락처 안내
        self._add_text_frame(
            slide, 0.5, 6.5, 9, 0.5,
            "상세 컨설팅 문의: contact@example.com",
            font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC),
            alignment=PP_ALIGN.CENTER
        )

    def generate(self, data: Dict[str, Any]) -> str:
        """
        PPT 생성

        Args:
            data: 진단 데이터

        Returns:
            생성된 PPT 파일 경로
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 9개 슬라이드 생성
        self._create_cover_slide(prs, data)            # 1. 표지 (충격 문구 포함)
        self._create_analysis_slide(prs, data)         # 2. 현황 분석 (7개 항목)
        self._create_detail_check_slide(prs, data)     # 3. 서비스 세팅 체크리스트
        self._create_competitor_slide(prs, data)       # 4. 경쟁사 비교 (신규)
        self._create_keyword_slide(prs, data)          # 5. 키워드 & 순위
        self._create_related_keywords_slide(prs, data) # 6. 연관 키워드 제안
        self._create_improvement_slide(prs, data)      # 7. 비어있는 항목 리스트
        self._create_breakeven_slide(prs, data)        # 8. 손익분기 계산 (신규)
        self._create_proposal_slide(prs, data)         # 9. 제안

        # 파일명 생성
        business_name = data.get("business_name", "unknown")
        # 파일명에 사용할 수 없는 문자 제거
        safe_name = "".join(c for c in business_name if c.isalnum() or c in (" ", "_", "-")).strip()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_name}_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)

        # 저장
        prs.save(filepath)

        return filename
