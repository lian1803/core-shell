"""
generators/ppt_generator.py — 11슬라이드 PPT 생성 (완전 재작성)
기존 템플릿 기반 + 동적 차트/색상 주입
"""
import sys
import os
import re
from typing import Dict, Any, List, Optional

sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import COLOR, GRADE_COLOR, SLIDE_INDEX, get_grade, log


# ── 텍스트 교체 헬퍼 ─────────────────────────────────────────

def _replace_in_shape(shape, replacements: Dict[str, str]):
    """shape 내부 텍스트 일괄 교체 (폰트 서식 유지, 그룹 재귀)"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            _replace_in_shape(child, replacements)


def replace_all_text(prs: Presentation, replacements: Dict[str, str]):
    """PPT 전체 슬라이드에서 텍스트 교체"""
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, replacements)


def _set_shape_color(shape, rgb: RGBColor):
    """도형 채우기 색상 변경"""
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    except Exception:
        pass


def _set_text_color(shape, rgb: RGBColor):
    """shape 내 모든 텍스트 색상 변경"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                run.font.color.rgb = rgb
            except Exception:
                pass


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=14, bold=False, color=None, bg_color=None):
    """텍스트박스 신규 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


# ── 슬라이드별 생성 함수 ────────────────────────────────────

def _build_slide_cover(slide, business_name: str, category: str):
    """슬라이드 1: 표지 — 업체명 교체"""
    replacements = {
        "(업체명)": business_name,
        "업체명": business_name,
        "한의원": category or "업종",
    }
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)


def _build_slide_score(slide, score_result, business_name: str, lost_customers: int = 0):
    """슬라이드 2: 종합 점수 + 등급 + 기회비용"""
    total = score_result.total_score
    grade = score_result.grade
    grade_color = score_result.grade_color
    relative_pct = score_result.relative_pct

    # 기존 텍스트 교체
    replacements = {
        "(업체명)": business_name,
        "업체명": business_name,
        "점수": f"{total}점",
        "등급": f"{grade} 등급",
    }
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)

    # 중앙 상단에 등급 박스
    _add_textbox(
        slide, f"{grade}",
        left=Inches(1), top=Inches(1.2),
        width=Inches(2), height=Inches(1.5),
        font_size=72, bold=True,
        color=COLOR["WHITE"], bg_color=grade_color,
    )
    # 점수 (실제 만점 95점)
    _add_textbox(
        slide, f"{total} / 95점",
        left=Inches(3.2), top=Inches(1.4),
        width=Inches(3.5), height=Inches(1.0),
        font_size=40, bold=True,
        color=COLOR["DARK"],
    )
    # 경쟁사 대비
    if relative_pct > 0:
        _add_textbox(
            slide, f"경쟁 상위 업체 평균 대비  {relative_pct}%  수준",
            left=Inches(1), top=Inches(2.8),
            width=Inches(8), height=Inches(0.6),
            font_size=16, bold=False,
            color=COLOR["BAD"] if relative_pct < 80 else COLOR["GOOD"],
        )

    # 기회비용 (잃고 있는 고객 수)
    if lost_customers > 0:
        _add_textbox(
            slide,
            f"⚠  현재 순위 기준, 매달 약 {lost_customers:,}명의 고객을 놓치고 있습니다",
            left=Inches(0.5), top=Inches(3.5),
            width=Inches(9), height=Inches(0.65),
            font_size=14, bold=True,
            color=COLOR["BAD"],
        )

    # 카테고리별 점수 바 (텍스트박스 너비 비율 표현)
    cat_scores = score_result.category_scores
    cat_max = {"기본정보": 20, "콘텐츠": 35, "운영관리": 20, "플랫폼연동": 12, "외부채널": 8}
    y_start = Inches(4.3)
    for cat_name, score in cat_scores.items():
        max_s = cat_max.get(cat_name, 20)
        pct_w = score / max_s if max_s > 0 else 0
        bar_max_w = Inches(5)
        bar_w = max(int(bar_max_w * pct_w), Inches(0.1))

        _add_textbox(
            slide, f"{cat_name}  {score}/{max_s}",
            left=Inches(1), top=y_start,
            width=Inches(3), height=Inches(0.4),
            font_size=11, color=COLOR["DARK"],
        )
        # 색상 바
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(4.2), y_start + Inches(0.05),
            bar_w, Inches(0.3),
        )
        bar_color = COLOR["GOOD"] if pct_w >= 0.7 else (COLOR["WARNING"] if pct_w >= 0.4 else COLOR["BAD"])
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        y_start += Inches(0.5)


def _build_slide_competitor(slide, place_data, competitors: list, score_result):
    """슬라이드 3: 경쟁사 비교 차트"""
    if not competitors:
        # 경쟁사 없으면 안내 텍스트만
        _add_textbox(
            slide, "경쟁사 데이터를 수집하지 못했습니다.\n(네트워크 오류 또는 검색 결과 없음)",
            left=Inches(1), top=Inches(2),
            width=Inches(8), height=Inches(1.5),
            font_size=16, color=COLOR["NEUTRAL"],
        )
        return

    comp_avg = score_result.competitor_avg

    my_review = getattr(place_data, 'visitor_review_count', 0)
    my_photo = getattr(place_data, 'photo_count', 0)
    my_save = getattr(place_data, 'save_count', 0)

    avg_review = comp_avg.get("visitor_review_count", 0)
    avg_photo = comp_avg.get("photo_count", 0)
    avg_save = comp_avg.get("save_count", 0)

    # 막대 차트 생성 (python-pptx 내장)
    try:
        chart_data = ChartData()
        chart_data.categories = ['방문자 리뷰', '사진 수', '저장 수']
        chart_data.add_series('우리 업체', (my_review, my_photo, my_save))
        chart_data.add_series('경쟁사 평균', (avg_review, avg_photo, avg_save))

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.2),
            Inches(6.5), Inches(4.5),
            chart_data,
        )
        chart = chart_shape.chart

        # 시리즈 색상 지정
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = COLOR["NEUTRAL"]
        chart.series[1].format.fill.solid()
        chart.series[1].format.fill.fore_color.rgb = COLOR["GOOD"]

        log("PPT", "경쟁사 비교 차트 삽입 완료")
    except Exception as e:
        log("PPT", f"차트 생성 실패 (텍스트 대체): {e}")
        _add_textbox(
            slide,
            f"우리 업체 vs 경쟁사 평균\n"
            f"방문자 리뷰: {my_review}개 vs {avg_review:.0f}개\n"
            f"사진 수: {my_photo}장 vs {avg_photo:.0f}장\n"
            f"저장 수: {my_save}개 vs {avg_save:.0f}개",
            left=Inches(0.5), top=Inches(1.5),
            width=Inches(6), height=Inches(3),
            font_size=14, color=COLOR["DARK"],
        )

    # 우측에 퍼센트 차이 텍스트박스
    def pct_diff(mine, avg):
        if avg <= 0:
            return "N/A"
        diff = int((mine - avg) / avg * 100)
        return f"{diff:+d}%"

    diff_lines = [
        f"리뷰: {pct_diff(my_review, avg_review)}",
        f"사진: {pct_diff(my_photo, avg_photo)}",
        f"저장: {pct_diff(my_save, avg_save)}",
    ]
    diff_text = "\n".join(diff_lines)
    _add_textbox(
        slide, diff_text,
        left=Inches(7.2), top=Inches(2.0),
        width=Inches(2.5), height=Inches(2.0),
        font_size=14, bold=True,
        color=COLOR["BAD"],
    )


def _build_slide_biz_info(slide, place_data, keyword_stats: Dict, main_keyword: str, rank: int):
    """슬라이드 4: 업체 기본정보 + 키워드 조회수"""
    name = getattr(place_data, 'name', '') or ''
    category = getattr(place_data, 'category', '') or ''
    address = getattr(place_data, 'address', '') or ''
    phone = getattr(place_data, 'phone', '') or ''

    pc = keyword_stats.get("pc", 0)
    mobile = keyword_stats.get("mobile", 0)
    total = keyword_stats.get("total", 0)
    rank_str = f"{rank}위" if rank > 0 else "순위권 외"

    replacements = {
        "(업체명)": name,
        "업체명": name,
        "한의원": category,
        "59위": rank_str,
        "Pc : 80회": f"PC : {pc:,}회",
        "모바일 : 80회": f"모바일 : {mobile:,}회",
        "전체 : 80회": f"전체 : {total:,}회",
        "플레이스 주소": address,
        "상호명": name,
    }
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)

    # 키워드 + 조회수 추가 표시
    _add_textbox(
        slide,
        f"메인 키워드: {main_keyword}\n"
        f"월 검색량: PC {pc:,}회 / 모바일 {mobile:,}회\n"
        f"현재 플레이스 순위: {rank_str}",
        left=Inches(0.5), top=Inches(5.2),
        width=Inches(9), height=Inches(1.2),
        font_size=12, color=COLOR["DARK"],
    )


def _build_slide_keywords(slide, related_keywords: List[Dict], business_name: str):
    """슬라이드 5: 확장 키워드 (기존 fill_slide3_keywords 로직 유지)"""
    keyword_boxes = []
    volume_boxes = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = " ".join(p.text for p in shape.text_frame.paragraphs).strip()
        if "디자인 시벌" in text:
            keyword_boxes.append(shape)
        elif text == "20,180" or (text.replace(",", "").isdigit() and 4 <= len(text.replace(",", "")) <= 6):
            volume_boxes.append(shape)

    for i, shape in enumerate(keyword_boxes):
        kw = related_keywords[i]["keyword"] if i < len(related_keywords) else "-"
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = kw
        if not any(run.text for para in shape.text_frame.paragraphs for run in para.runs):
            shape.text_frame.paragraphs[0].text = kw

    for i, shape in enumerate(volume_boxes):
        vol = f"{related_keywords[i]['mobile']:,}" if i < len(related_keywords) else "-"
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = vol
        if not any(run.text for para in shape.text_frame.paragraphs for run in para.runs):
            shape.text_frame.paragraphs[0].text = vol

    # 키워드 부족 시 나머지 채우기
    if not keyword_boxes and related_keywords:
        lines = [f"{k['keyword']}  |  {k['mobile']:,}회" for k in related_keywords[:7]]
        _add_textbox(
            slide, "\n".join(lines),
            left=Inches(1), top=Inches(1.5),
            width=Inches(8), height=Inches(4),
            font_size=14, color=COLOR["DARK"],
        )


def _build_slide_screenshot(slide, screenshot_path: str, main_keyword: str):
    """슬라이드 6: 검색 결과 스크린샷"""
    if screenshot_path and os.path.exists(screenshot_path):
        try:
            slide.shapes.add_picture(
                screenshot_path,
                Inches(0.3), Inches(1.0),
                Inches(4.5), Inches(5.5),
            )
            log("PPT", "스크린샷 삽입 완료")
        except Exception as e:
            log("PPT", f"스크린샷 삽입 실패: {e}")
    else:
        _add_textbox(
            slide,
            f"'{main_keyword}' 검색 결과\n(스크린샷 수집 실패)",
            left=Inches(1), top=Inches(2),
            width=Inches(5), height=Inches(2),
            font_size=14, color=COLOR["NEUTRAL"],
        )


def _build_slide_checklist(slide, score_result, place_data):
    """슬라이드 7: 22항목 진단 체크리스트 (O/X + 색상)"""
    item_scores = score_result.item_scores
    item_max = score_result.item_max

    # 한글 항목명 → 실제 데이터 매핑
    check_map = {
        "업체명":     True,
        "카테고리":   bool(getattr(place_data, 'category', '')),
        "주소":       bool(getattr(place_data, 'address', '')),
        "전화번호":   bool(getattr(place_data, 'phone', '')),
        "영업시간":   getattr(place_data, 'has_hours', False),
        "메뉴_가격":  getattr(place_data, 'has_menu', False),
        "소개글":     getattr(place_data, 'has_intro', False),
        "사진수":     getattr(place_data, 'photo_count', 0) >= 5,
        "방문자리뷰": getattr(place_data, 'visitor_review_count', 0) >= 10,
        "블로그리뷰": getattr(place_data, 'blog_review_count', 0) >= 5,
        "영수증리뷰": getattr(place_data, 'receipt_review_count', 0) >= 10,
        "저장수":     getattr(place_data, 'save_count', 0) >= 10,
        "새소식":     getattr(place_data, 'news_last_days', 9999) <= 90,
        "사장님답글": getattr(place_data, 'owner_reply_rate', 0.0) >= 0.2,
        "해시태그":   getattr(place_data, 'has_hashtag', False),
        "키워드등록": len(getattr(place_data, 'keywords', [])) >= 1,
        "네이버예약": getattr(place_data, 'has_reservation', False),
        "네이버톡톡": getattr(place_data, 'has_talktalk', False),
        "스마트콜":   getattr(place_data, 'has_smartcall', False),
        "쿠폰이벤트": getattr(place_data, 'has_coupon', False),
        "인스타그램": getattr(place_data, 'has_instagram', False),
        "카카오채널": getattr(place_data, 'has_kakao', False),
    }

    # 카테고리 그룹
    groups = {
        "기본 정보 (20점)": ["업체명","카테고리","주소","전화번호","영업시간","메뉴_가격","소개글"],
        "콘텐츠 활성도 (35점)": ["사진수","방문자리뷰","블로그리뷰","영수증리뷰","저장수","새소식"],
        "관리 운영 (20점)": ["사장님답글","해시태그","키워드등록"],
        "플랫폼 연동 (20점)": ["네이버예약","네이버톡톡","스마트콜","쿠폰이벤트","인스타그램","카카오채널"],
    }

    y = Inches(1.0)
    col1_x = Inches(0.3)
    col2_x = Inches(5.0)
    row_h = Inches(0.32)

    for group_name, items in groups.items():
        # 그룹 헤더
        _add_textbox(
            slide, group_name,
            left=col1_x, top=y,
            width=Inches(9), height=Inches(0.35),
            font_size=11, bold=True,
            color=COLOR["DARK"],
        )
        y += Inches(0.38)

        # 항목 2열 배치
        col = 0
        for item in items:
            passed = check_map.get(item, False)
            mark = "O" if passed else "X"
            mark_color = COLOR["GOOD"] if passed else COLOR["BAD"]
            score_v = item_scores.get(item, 0)
            max_v = item_max.get(item, 0)
            item_label = item.replace("_", "/")
            txt = f"[{mark}] {item_label}  ({score_v}/{max_v}점)"

            x_pos = col1_x if col == 0 else col2_x
            _add_textbox(
                slide, txt,
                left=x_pos, top=y,
                width=Inches(4.5), height=row_h,
                font_size=10,
                color=mark_color,
            )
            col += 1
            if col >= 2:
                col = 0
                y += row_h + Inches(0.02)

        if col != 0:
            y += row_h + Inches(0.02)
        y += Inches(0.15)


def _build_slide_weak_points(slide, score_result, place_data, comp_avg: Dict):
    """슬라이드 8: 치명적 3가지 취약 항목"""
    weak = score_result.weak_points[:3]
    if not weak:
        weak = ["사진수", "방문자리뷰", "소개글"]

    # 취약 항목 → 현재 상태 텍스트 매핑
    def _desc(item_name: str) -> str:
        mapping = {
            "사진수":     f"현재 {getattr(place_data,'photo_count',0)}장",
            "방문자리뷰": f"현재 {getattr(place_data,'visitor_review_count',0)}개",
            "블로그리뷰": f"현재 {getattr(place_data,'blog_review_count',0)}개",
            "영수증리뷰": f"현재 {getattr(place_data,'receipt_review_count',0)}개",
            "저장수":     f"현재 {getattr(place_data,'save_count',0)}개",
            "해시태그":   "미설정",
            "새소식":     f"{getattr(place_data,'news_last_days',9999)}일 전 업데이트",
            "사장님답글": f"응답률 {int(getattr(place_data,'owner_reply_rate',0)*100)}%",
            "네이버예약": "미연동",
            "인스타그램": "미연동",
            "카카오채널": "미연동",
            "소개글":     "미등록" if not getattr(place_data,'has_intro',False) else "등록됨",
            "메뉴_가격":  "미등록" if not getattr(place_data,'has_menu',False) else "등록됨",
        }
        return mapping.get(item_name, "미완료")

    y = Inches(1.3)
    for i, item in enumerate(weak):
        score_v = score_result.item_scores.get(item, 0)
        max_v = score_result.item_max.get(item, 0)
        desc = _desc(item)
        item_label = item.replace("_", "/")

        # 번호 + 항목명
        _add_textbox(
            slide, f"  {i+1}. {item_label}",
            left=Inches(0.5), top=y,
            width=Inches(3.5), height=Inches(0.5),
            font_size=18, bold=True,
            color=COLOR["BAD"],
        )
        # 현재 상태
        _add_textbox(
            slide, f"현재: {desc}  |  점수: {score_v}/{max_v}점",
            left=Inches(4.2), top=y + Inches(0.08),
            width=Inches(5), height=Inches(0.4),
            font_size=13,
            color=COLOR["DARK"],
        )
        # 경쟁사 비교 (해당 항목이 있는 경우)
        comp_val_map = {
            "사진수":     ("photo_count", "장"),
            "방문자리뷰": ("visitor_review_count", "개"),
            "저장수":     ("save_count", "개"),
            "블로그리뷰": ("blog_review_count", "개"),
        }
        if item in comp_val_map and comp_avg:
            field_key, unit = comp_val_map[item]
            avg_val = comp_avg.get(field_key, 0)
            if avg_val > 0:
                _add_textbox(
                    slide, f"→ 경쟁사 평균 {avg_val:.0f}{unit} 수준",
                    left=Inches(4.2), top=y + Inches(0.52),
                    width=Inches(5), height=Inches(0.35),
                    font_size=12, color=COLOR["WARNING"],
                )

        y += Inches(1.6)


def _build_slide_package(slide, business_name: str):
    """슬라이드 9: 패키지 소개"""
    replacements = {"(업체명)": business_name, "업체명": business_name}
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)

    # 슬라이드가 비어있으면 패키지 콘텐츠 직접 생성
    has_content = any(
        shape.has_text_frame and any(p.text.strip() for p in shape.text_frame.paragraphs)
        for shape in slide.shapes
    )
    if not has_content:
        _add_textbox(slide, "서비스 패키지", Inches(0.5), Inches(0.2),
                     Inches(9), Inches(0.7), font_size=28, bold=True, color=COLOR["DARK"])
        _add_textbox(slide, f"[ {business_name} ] 맞춤 플레이스 최적화 플랜",
                     Inches(0.5), Inches(0.95), Inches(9), Inches(0.45),
                     font_size=13, color=COLOR["NEUTRAL"])

        pkgs = [
            ("주목 패키지  290,000원/월",
             "• 기본 정보 전체 최적화  • 키워드 3개 등록\n"
             "• 사진 정리 및 순서 최적화  • 월 1회 리포트 제공"),
            ("집중 패키지  490,000원/월",
             "• 주목 패키지 전체 포함  • 전문 사진 촬영 10장\n"
             "• 블로그 리뷰 5건 관리  • 사장님 답글 대행"),
            ("시선 패키지  890,000원/월",
             "• 집중 패키지 전체 포함  • 인스타그램 연동 관리\n"
             "• 월 2회 새소식/이벤트 게시  • 전담 매니저 배정"),
        ]
        y = Inches(1.5)
        for title, detail in pkgs:
            _add_textbox(slide, title, Inches(0.5), y, Inches(9), Inches(0.5),
                         font_size=16, bold=True, color=COLOR["PRIMARY"])
            _add_textbox(slide, detail, Inches(0.9), y + Inches(0.52), Inches(8), Inches(0.75),
                         font_size=12, color=COLOR["DARK"])
            y += Inches(1.55)


def _build_slide_package_detail(slide, business_name: str):
    """슬라이드 10: 패키지 도입 근거"""
    replacements = {"(업체명)": business_name, "업체명": business_name}
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)

    has_content = any(
        shape.has_text_frame and any(p.text.strip() for p in shape.text_frame.paragraphs)
        for shape in slide.shapes
    )
    if not has_content:
        _add_textbox(slide, "왜 네이버 플레이스 최적화인가?",
                     Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
                     font_size=26, bold=True, color=COLOR["DARK"])

        points = [
            ("01", "네이버 지도 검색 = 구매 의도가 가장 높은 고객",
             "검색하는 사람 10명 중 7명은 '오늘 방문' 또는 '바로 구매' 의사가 있습니다."),
            ("02", "상위 3위 이내 노출 시 클릭률 35% 이상",
             "4위 이하로 떨어지는 순간 클릭률이 절반 이하로 감소합니다."),
            ("03", "광고비 없는 자연 유입 채널",
             "한 번 최적화된 플레이스는 지속적으로 고객을 유입시킵니다."),
            ("04", "경쟁사는 이미 시작했습니다",
             "상위 업체들은 대부분 전문 관리 중입니다. 지금이 마지막 기회입니다."),
        ]
        y = Inches(1.1)
        for num, title, desc in points:
            _add_textbox(slide, num, Inches(0.3), y, Inches(0.7), Inches(0.45),
                         font_size=20, bold=True, color=COLOR["PRIMARY"])
            _add_textbox(slide, title, Inches(1.1), y, Inches(8.3), Inches(0.42),
                         font_size=14, bold=True, color=COLOR["DARK"])
            _add_textbox(slide, desc, Inches(1.1), y + Inches(0.44), Inches(8.3), Inches(0.38),
                         font_size=11, color=COLOR["NEUTRAL"])
            y += Inches(1.1)


def _build_slide_estimate(slide, business_name: str, grade: str, lost_customers: int = 0):
    """슬라이드 11: 견적서 + 기회비용 요약"""
    pkg_map = {"S": "시선 패키지", "A": "집중 패키지", "B": "집중 패키지",
               "C": "주목 패키지", "D": "주목 패키지", "F": "주목 패키지"}
    price_map = {"시선 패키지": "890,000원/월", "집중 패키지": "490,000원/월",
                 "주목 패키지": "290,000원/월"}
    pkg = pkg_map.get(grade, "주목 패키지")
    price = price_map[pkg]

    replacements = {
        "(업체명)": business_name,
        "업체명": business_name,
        "집중패키지": pkg,
        "패키지명": pkg,
    }
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)

    has_content = any(
        shape.has_text_frame and any(p.text.strip() for p in shape.text_frame.paragraphs)
        for shape in slide.shapes
    )
    if not has_content:
        _add_textbox(slide, "맞춤 제안서",
                     Inches(0.5), Inches(0.2), Inches(9), Inches(0.65),
                     font_size=28, bold=True, color=COLOR["DARK"])
        _add_textbox(slide, f"[ {business_name} ] 전용 제안",
                     Inches(0.5), Inches(0.9), Inches(9), Inches(0.4),
                     font_size=13, color=COLOR["NEUTRAL"])

        # 추천 패키지 박스
        _add_textbox(slide, f"추천 패키지: {pkg}",
                     Inches(0.5), Inches(1.45), Inches(6), Inches(0.65),
                     font_size=22, bold=True, color=COLOR["WHITE"],
                     bg_color=COLOR["PRIMARY"])
        _add_textbox(slide, price,
                     Inches(6.7), Inches(1.52), Inches(2.7), Inches(0.52),
                     font_size=18, bold=True, color=COLOR["PRIMARY"])

        # 기회비용 계산
        if lost_customers > 0:
            _add_textbox(
                slide,
                f"지금 계약하지 않으면,\n매달 약 {lost_customers:,}명의 잠재 고객이 경쟁사로 갑니다.",
                Inches(0.5), Inches(2.25), Inches(9), Inches(0.9),
                font_size=14, bold=True, color=COLOR["BAD"],
            )
            y_after = Inches(3.3)
        else:
            y_after = Inches(2.35)

        # 서비스 포함 내역
        _add_textbox(slide, "포함 서비스", Inches(0.5), y_after, Inches(9), Inches(0.42),
                     font_size=14, bold=True, color=COLOR["DARK"])
        include_map = {
            "시선 패키지": "기본정보 최적화 / 사진 촬영 10장 / 블로그 리뷰 5건 / 인스타그램 관리 / 월 2회 새소식 / 전담 매니저",
            "집중 패키지": "기본정보 최적화 / 사진 촬영 10장 / 블로그 리뷰 5건 / 사장님 답글 대행 / 월 1회 리포트",
            "주목 패키지": "기본정보 최적화 / 키워드 3개 등록 / 사진 정리 / 월 1회 리포트",
        }
        _add_textbox(slide, include_map[pkg],
                     Inches(0.7), y_after + Inches(0.45), Inches(8.5), Inches(0.5),
                     font_size=12, color=COLOR["DARK"])

    # 기존 템플릿이 있는 경우도 추천 패키지 강조 추가
    _add_textbox(
        slide,
        f"[ {business_name} ] 추천: {pkg}  ({price})",
        left=Inches(0.5), top=Inches(6.8),
        width=Inches(9), height=Inches(0.5),
        font_size=13, bold=True,
        color=COLOR["PRIMARY"],
    )


# ── 메인 PPT 생성 함수 ─────────────────────────────────────

def generate_ppt(
    template_path: str,
    output_path: str,
    business_name: str,
    category: str,
    place_data,
    keyword_stats: Dict,
    main_keyword: str,
    rank: int,
    score_result,
    competitors: list,
    screenshot_path: str = "",
) -> str:
    """
    11슬라이드 PPT 생성.
    템플릿에 슬라이드가 11개 미만이면 기존 슬라이드만 처리.
    반환값: 저장된 파일 경로
    """
    try:
        prs = Presentation(template_path)
        slide_count = len(prs.slides)
        log("PPT", f"템플릿 슬라이드 수: {slide_count}")

        # CRITICAL 검증: 예상 슬라이드 수 vs 실제 템플릿 슬라이드 수
        expected_count = max(SLIDE_INDEX.values()) + 1  # 11
        if slide_count < expected_count:
            missing_slides = [name for name, idx in SLIDE_INDEX.items() if idx >= slide_count]
            log("PPT", f"⚠️ [경고] 템플릿 슬라이드 {slide_count}장 (예상 {expected_count}장)")
            log("PPT", f"⚠️ 누락된 슬라이드 {len(missing_slides)}개: {missing_slides}")
            log("PPT", "   → 해당 슬라이드는 건너뜁니다. 템플릿에 슬라이드를 추가하면 생성됩니다.")

        # 먼저 전체 PPT에서 공통 플레이스홀더 일괄 교체
        global_replacements = {
            "(업체명)": business_name,
            "업체명": business_name,
        }
        replace_all_text(prs, global_replacements)
        log("PPT", f"전체 슬라이드에서 플레이스홀더 교체 완료: {business_name}")

        related = keyword_stats.get("related", [])
        # 연관 키워드 부족 시 place 키워드로 보완
        if len(related) < 4 and getattr(place_data, 'keywords', []):
            extra = [{"keyword": k, "mobile": 0} for k in place_data.keywords[:7]]
            related = (related + extra)[:7]

        # 슬라이드별 처리 (인덱스 존재 확인 후)
        def _has_slide(idx: int) -> bool:
            return idx < slide_count

        # 슬라이드 1: 표지
        if _has_slide(SLIDE_INDEX["cover"]):
            _build_slide_cover(prs.slides[SLIDE_INDEX["cover"]], business_name, category)

        # 슬라이드 2: 종합 점수
        if _has_slide(SLIDE_INDEX["score"]):
            _build_slide_score(prs.slides[SLIDE_INDEX["score"]], score_result, business_name)

        # 슬라이드 3: 경쟁사 비교 차트
        if _has_slide(SLIDE_INDEX["competitor"]):
            _build_slide_competitor(
                prs.slides[SLIDE_INDEX["competitor"]],
                place_data, competitors, score_result
            )

        # 슬라이드 4: 업체 기본정보 + 키워드
        if _has_slide(SLIDE_INDEX["biz_info"]):
            _build_slide_biz_info(
                prs.slides[SLIDE_INDEX["biz_info"]],
                place_data, keyword_stats, main_keyword, rank
            )

        # 슬라이드 5: 확장 키워드
        if _has_slide(SLIDE_INDEX["keywords"]):
            _build_slide_keywords(
                prs.slides[SLIDE_INDEX["keywords"]],
                related, business_name
            )

        # 슬라이드 6: 검색 결과 스크린샷
        if _has_slide(SLIDE_INDEX["screenshot"]):
            _build_slide_screenshot(
                prs.slides[SLIDE_INDEX["screenshot"]],
                screenshot_path, main_keyword
            )

        # 슬라이드 7: 22항목 체크리스트
        if _has_slide(SLIDE_INDEX["checklist"]):
            _build_slide_checklist(
                prs.slides[SLIDE_INDEX["checklist"]],
                score_result, place_data
            )

        # 슬라이드 8: 치명적 3가지
        if _has_slide(SLIDE_INDEX["weak_points"]):
            _build_slide_weak_points(
                prs.slides[SLIDE_INDEX["weak_points"]],
                score_result, place_data,
                score_result.competitor_avg
            )

        # 슬라이드 9~11: 패키지 / 견적서
        if _has_slide(SLIDE_INDEX["package"]):
            _build_slide_package(prs.slides[SLIDE_INDEX["package"]], business_name)

        if _has_slide(SLIDE_INDEX["package_detail"]):
            _build_slide_package_detail(prs.slides[SLIDE_INDEX["package_detail"]], business_name)

        if _has_slide(SLIDE_INDEX["estimate"]):
            _build_slide_estimate(
                prs.slides[SLIDE_INDEX["estimate"]],
                business_name, score_result.grade
            )

        prs.save(output_path)
        log("PPT", f"저장 완료: {output_path}")
        return output_path

    except Exception as e:
        log("PPT", f"생성 오류: {e}")
        raise
