"""
네이버 플레이스 URL → 자동 분석 → PPT 제안서 생성
사용법: python process.py
"""
import asyncio
import sys
import os
import re
import time
import hmac
import hashlib
import base64
import random
import urllib.parse
from copy import deepcopy
from datetime import datetime
from typing import Optional, Dict, Any, List

sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import httpx
import openpyxl
from pptx import Presentation
from pptx.util import Pt
from playwright.async_api import async_playwright

# ──────────────────────────────────────────────
# 경로 설정
# ──────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
EXCEL_PATH = os.path.join(BASE_DIR, "샘플", "샘플.xlsx")
TEMPLATE_PATH = os.path.join(BASE_DIR, "샘플", "Black and Grey Minimalist Company Project Proposal.pptx")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Naver API 키 (naver-diagnosis .env에서 가져옴)
NAVER_CLIENT_ID = "o0776HJDmQVO6J9Lez1m"
NAVER_CLIENT_SECRET = "ZXG0lPbgH9"
NAVER_AD_CUSTOMER_ID = "3280509"
NAVER_AD_ACCESS_LICENSE = "0100000000265e400df7ebfb6855d1d16f25948785a08c487285d1876915996b92c0a433fd"
NAVER_AD_SECRET_KEY = "AQAAAAAmXkAN9+v7aFXR0W8llIeF0zQ4nUlBHoEFBwNqw1pf6A=="

USER_AGENTS = [
    "Mozilla/5.0 (iPhone; CPU iPhone OS 17_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.0 Mobile/15E148 Safari/604.1",
    "Mozilla/5.0 (Linux; Android 14; SM-S918B) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0.0.0 Mobile Safari/537.36",
]


# ──────────────────────────────────────────────
# 1. 엑셀에서 첫 번째 URL 있는 행 읽기
# ──────────────────────────────────────────────
def read_first_url_row(path: str) -> Optional[Dict]:
    wb = openpyxl.load_workbook(path)
    ws = wb.active
    headers = [str(c.value) for c in ws[1]]
    print(f"[Excel] 컬럼: {headers}")

    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = dict(zip(headers, row))
        # 모든 컬럼에서 네이버 URL 탐색
        url = ""
        for col_name in ["네이버플레이스", "주소", "플레이스주소", "url", "URL"]:
            val = str(row_dict.get(col_name, "") or "")
            if val.startswith("http") and "naver" in val:
                url = val
                break
        # 컬럼명 상관없이 값에서 URL 탐색
        if not url:
            for val in row_dict.values():
                val_str = str(val or "")
                if val_str.startswith("http") and "naver" in val_str:
                    url = val_str
                    break
        if url and "place" in url:
            row_dict["_place_url"] = url
            print(f"[Excel] 발견: {row_dict.get('업체명')} | URL: {url[:70]}...")
            return row_dict
    return None


# ──────────────────────────────────────────────
# 2. URL에서 place_id 추출
# ──────────────────────────────────────────────
def extract_place_id(url: str) -> Optional[str]:
    match = re.search(r"place/(\d+)", url)
    return match.group(1) if match else None


# ──────────────────────────────────────────────
# 3. 네이버 플레이스 크롤링
# ──────────────────────────────────────────────
async def crawl_place(browser, place_id: str) -> Dict[str, Any]:
    url = f"https://m.place.naver.com/place/{place_id}/home"
    result = {
        "place_id": place_id,
        "name": "",
        "category": "",
        "address": "",
        "phone": "",
        "photo_count": 0,
        "visitor_review_count": 0,
        "blog_review_count": 0,
        "receipt_review_count": 0,
        "keywords": [],
        "has_hours": False,
        "has_menu": False,
        "has_intro": False,
        "has_booking": False,
        "has_talktalk": False,
        "has_smartcall": False,
        "has_coupon": False,
        "has_news": False,
        "has_owner_reply": False,
        "has_instagram": False,
    }
    page = None
    try:
        context = await browser.new_context(
            user_agent=random.choice(USER_AGENTS),
            viewport={"width": 390, "height": 844},
            locale="ko-KR",
        )
        page = await context.new_page()
        print(f"[Crawl] {url}")
        await page.goto(url, timeout=30000)
        await page.wait_for_load_state("networkidle", timeout=30000)

        text = await page.inner_text("body")
        content = await page.content()

        # JSON에서 업체 기본정보 추출 (name, address, phone, category)
        import json
        json_match = re.search(r'"name"\s*:\s*"([^"]+)"', content)
        if json_match:
            result["name"] = json_match.group(1)

        addr_match = re.search(r'"roadAddress"\s*:\s*"([^"]+)"', content)
        if not addr_match:
            addr_match = re.search(r'"address"\s*:\s*"([^"]+)"', content)
        if addr_match:
            result["address"] = addr_match.group(1)

        phone_match = re.search(r'"phone"\s*:\s*"([0-9\-]+)"', content)
        if phone_match:
            result["phone"] = phone_match.group(1)

        cat_match = re.search(r'"category"\s*:\s*"([^"]+)"', content)
        if cat_match:
            result["category"] = cat_match.group(1)

        # 사진 수: /photo 페이지에서 SasImage total 패턴으로 추출
        for photo_url in [
            f"https://m.place.naver.com/restaurant/{place_id}/photo",
            f"https://m.place.naver.com/place/{place_id}/photo",
        ]:
            try:
                p2 = await context.new_page()
                await p2.goto(photo_url, timeout=30000)
                await p2.wait_for_load_state("networkidle", timeout=30000)
                photo_html = await p2.content()
                await p2.close()

                m1 = re.findall(r'SasImage[^}]*total["\s:]+(\d+)', photo_html)
                if m1:
                    result["photo_count"] = max(int(v) for v in m1)
                    break
                m2 = re.findall(r'"relation"\s*:\s*"[^"]*사진[^"]*"[^}]*"total"\s*:\s*(\d+)', photo_html)
                if m2:
                    result["photo_count"] = max(int(v) for v in m2)
                    break
            except Exception:
                pass

        # 리뷰 수
        rv = re.search(r"방문자 리뷰\s*([\d,]+)", text)
        if rv:
            result["visitor_review_count"] = int(rv.group(1).replace(",", ""))
        rv2 = re.search(r"블로그 리뷰\s*([\d,]+)", text)
        if rv2:
            result["blog_review_count"] = int(rv2.group(1).replace(",", ""))
        rv3 = re.search(r"영수증 리뷰\s*([\d,]+)", text)
        if rv3:
            result["receipt_review_count"] = int(rv3.group(1).replace(",", ""))

        # 영업시간
        result["has_hours"] = any(k in text for k in ["영업", "운영시간", "오전", "오후"])

        # 메뉴
        result["has_menu"] = "메뉴" in text

        # 소개글
        result["has_intro"] = any(k in text for k in ["소개", "안내", "정보"])

        # 예약/톡톡/스마트콜
        result["has_booking"] = "예약" in text
        result["has_talktalk"] = "톡톡" in text
        result["has_smartcall"] = "스마트콜" in text or "전화" in text

        # 쿠폰/새소식
        result["has_coupon"] = any(k in text for k in ["쿠폰", "이벤트", "할인"])
        result["has_news"] = "새소식" in text or "소식" in text

        # 사장님 답글
        result["has_owner_reply"] = "사장님" in text and ("답글" in text or "답변" in text)

        # 인스타그램
        result["has_instagram"] = "instagram" in content.lower() or "인스타" in text

        # 키워드 추출 (iframe 포함)
        frames_html = content
        try:
            for frame in page.frames:
                try:
                    fc = await frame.content()
                    frames_html += fc
                except Exception:
                    pass
        except Exception:
            pass

        kw_matches = re.findall(r'"keywordList"\s*:\s*\[([^\]]+)\]', frames_html)
        if kw_matches:
            keywords = re.findall(r'"([^"]{2,20})"', kw_matches[0])
            result["keywords"] = keywords[:10]

        if not result["keywords"]:
            # 데스크톱 버전에서 키워드 추출 시도
            try:
                dk_url = f"https://place.map.naver.com/place/{place_id}/home"
                p3 = await context.new_page()
                await p3.goto(dk_url, timeout=20000)
                await p3.wait_for_load_state("networkidle", timeout=15000)
                dk_html = await p3.content()
                for frame in p3.frames:
                    try:
                        dk_html += await frame.content()
                    except Exception:
                        pass
                kw2 = re.findall(r'"keywordList"\s*:\s*\[([^\]]+)\]', dk_html)
                if kw2:
                    keywords2 = re.findall(r'"([^"]{2,20})"', kw2[0])
                    result["keywords"] = keywords2[:10]
                await p3.close()
            except Exception:
                pass

        print(f"[Crawl] 완료: 사진={result['photo_count']}, 방문자리뷰={result['visitor_review_count']}, 블로그={result['blog_review_count']}, 키워드={result['keywords'][:3]}")

    except Exception as e:
        print(f"[Crawl] 오류: {e}")
    finally:
        if page:
            await page.close()
    return result


# ──────────────────────────────────────────────
# 4. 키워드 순위 조회
# ──────────────────────────────────────────────
async def get_keyword_rank(browser, keyword: str, business_name: str) -> int:
    """모바일 네이버 플레이스 검색에서 업체 순위 반환 (없으면 0)"""
    page = None
    try:
        context = await browser.new_context(
            user_agent=random.choice(USER_AGENTS),
            viewport={"width": 390, "height": 844},
            locale="ko-KR",
        )
        page = await context.new_page()
        encoded = urllib.parse.quote(keyword)
        url = f"https://m.search.naver.com/search.naver?query={encoded}&where=m_local"
        await page.goto(url, timeout=30000)
        await page.wait_for_load_state("networkidle", timeout=20000)

        text = await page.inner_text("body")
        lines = [l.strip() for l in text.split("\n") if l.strip()]

        # 업체명이 나오는 위치 찾기
        short_name = business_name.split()[0] if " " in business_name else business_name
        for i, line in enumerate(lines):
            if short_name in line:
                # 위치를 순위로 환산 (대략)
                print(f"[Rank] '{keyword}' 검색 → '{business_name}' 발견 (줄 {i})")
                # 플레이스 목록에서 몇 번째인지 계산
                place_items = [l for l in lines[:i] if len(l) > 2]
                return min(len(place_items) // 3 + 1, 99)

        print(f"[Rank] '{keyword}'에서 '{business_name}' 미발견")
        return 0
    except Exception as e:
        print(f"[Rank] 오류: {e}")
        return 0
    finally:
        if page:
            await page.close()


# ──────────────────────────────────────────────
# 5. 네이버 검색광고 API — 키워드 통계
# ──────────────────────────────────────────────
def _make_ad_headers(method: str, path: str) -> Dict[str, str]:
    timestamp = str(int(time.time() * 1000))
    message = f"{timestamp}.{method}.{path}"
    secret_bytes = NAVER_AD_SECRET_KEY.encode("utf-8")
    message_bytes = message.encode("utf-8")
    sig = hmac.new(secret_bytes, message_bytes, hashlib.sha256).digest()
    signature = base64.b64encode(sig).decode("utf-8")
    return {
        "Content-Type": "application/json; charset=UTF-8",
        "X-Timestamp": timestamp,
        "X-API-KEY": NAVER_AD_ACCESS_LICENSE,
        "X-Customer": NAVER_AD_CUSTOMER_ID,
        "X-Signature": signature,
    }


async def get_keyword_stats(keyword: str) -> Dict[str, Any]:
    path = "/keywordstool"
    params = {"hintKeywords": keyword, "showDetail": "1"}
    headers = _make_ad_headers("GET", path)
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                f"https://api.naver.com{path}",
                headers=headers,
                params=params,
                timeout=10.0,
            )
            data = resp.json()
        keywords_list = data.get("keywordList", [])
        if not keywords_list:
            return {"pc": 0, "mobile": 0, "total": 0, "related": []}

        # 정확히 일치하는 키워드 우선
        main = next((k for k in keywords_list if k.get("relKeyword") == keyword), keywords_list[0])
        pc = int(main.get("monthlyPcQcCnt", 0)) if str(main.get("monthlyPcQcCnt", "0")) != "< 10" else 5
        mobile = int(main.get("monthlyMobileQcCnt", 0)) if str(main.get("monthlyMobileQcCnt", "0")) != "< 10" else 5

        # 연관 키워드 (검색량 포함)
        related = []
        for k in keywords_list[:10]:
            kw_name = k.get("relKeyword", "")
            kw_mobile = k.get("monthlyMobileQcCnt", 0)
            if kw_mobile == "< 10":
                kw_mobile = 5
            else:
                try:
                    kw_mobile = int(kw_mobile)
                except Exception:
                    kw_mobile = 0
            if kw_name and kw_mobile > 0:
                related.append({"keyword": kw_name, "mobile": kw_mobile})

        related.sort(key=lambda x: x["mobile"], reverse=True)
        print(f"[KeywordAPI] '{keyword}': PC={pc}, 모바일={mobile}, 연관={len(related)}개")
        return {"pc": pc, "mobile": mobile, "total": pc + mobile, "related": related[:7]}
    except Exception as e:
        print(f"[KeywordAPI] 오류: {e}")
        return {"pc": 0, "mobile": 0, "total": 0, "related": []}


# ──────────────────────────────────────────────
# 6. PPT 텍스트 교체 헬퍼
# ──────────────────────────────────────────────
def replace_text_in_shape(shape, old: str, new: str):
    """shape 내 텍스트에서 old → new 교체 (폰트 서식 유지)"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def replace_all_text(prs: Presentation, replacements: Dict[str, str]):
    """PPT 전체에서 텍스트 교체"""
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, replacements)


def _replace_in_shape(shape, replacements: Dict[str, str]):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
    # Group 안쪽 shape들도 처리
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            _replace_in_shape(child, replacements)


def update_slide5_diagnosis(slide, data: Dict[str, Any]):
    """슬라이드5 종합진단 — 데이터 기반 텍스트 업데이트"""
    photo_count = data.get("photo_count", 0)
    has_intro = data.get("has_intro", False)
    has_coupon = data.get("has_coupon", False)
    has_news = data.get("has_news", False)

    # 대표사진 진단
    if photo_count >= 5:
        photo_msg = f"현재 {photo_count}장의 사진이 등록되어 있습니다.\n대표 사진을 마케팅 요소가 담긴 이미지로 최적화하면 더욱 효과적입니다."
    else:
        photo_msg = f"현재 {photo_count}장의 사진만 등록되어 있습니다.\n대표 사진 5장 이상 등록이 필요하며, 마케팅 요소가 있는 이미지로 재설정을 권장합니다."

    # 알림/이벤트 진단
    if has_coupon or has_news:
        event_msg = "알림 및 이벤트가 활성화되어 있습니다.\n지속적인 업데이트로 고객 재방문을 유도하세요."
    else:
        event_msg = "현재 알림 설정이 안 되어 있습니다.\n네이버에서 활용 가능한 플랫폼(알림, 이벤트)을 적극적으로\n활용하는 것이 좋습니다."

    # 소개/오시는 길 진단
    if has_intro:
        intro_msg = "소개 글이 등록되어 있습니다.\n대중교통 및 주차 정보, 키워드 최적화를 통해 더욱 강화하세요."
    else:
        intro_msg = "소개 글(정보란) 수정이 필요합니다.\n오시는 길(대중교통, 차량, 도보)에 대한\n정보와 키워드 재설정이 필요합니다."

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
        if "방문결정에 중요하게" in full_text or "마케팅적 요소가 들어간" in full_text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = photo_msg
        elif "현재 알림 설정" in full_text or "네이버에서 활용 가능한" in full_text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = event_msg
        elif "소개 글(정보란)" in full_text or "오시는 길(대중교통" in full_text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = intro_msg


def fill_slide3_keywords(slide, related_keywords: List[Dict], business_name: str):
    """슬라이드3 확장 키워드 채우기"""
    # '디자인 시벌' 텍스트박스들 (7개) → 키워드명으로
    # '20,180' 텍스트박스들 (7개) → 검색량으로
    keyword_boxes = []
    volume_boxes = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = " ".join(p.text for p in shape.text_frame.paragraphs).strip()
        if "디자인 시벌" in text:
            keyword_boxes.append(shape)
        elif text == "20,180" or (text.replace(",", "").isdigit() and 4 <= len(text.replace(",", "")) <= 6):
            volume_boxes.append(shape)

    # 키워드 채우기
    for i, shape in enumerate(keyword_boxes):
        if i < len(related_keywords):
            kw = related_keywords[i]["keyword"]
        else:
            kw = "-"
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = kw
        if not any(run.text for para in shape.text_frame.paragraphs for run in para.runs):
            shape.text_frame.paragraphs[0].text = kw

    # 검색량 채우기
    for i, shape in enumerate(volume_boxes):
        if i < len(related_keywords):
            vol = f"{related_keywords[i]['mobile']:,}"
        else:
            vol = "-"
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = vol
        if not any(run.text for para in shape.text_frame.paragraphs for run in para.runs):
            shape.text_frame.paragraphs[0].text = vol


# ──────────────────────────────────────────────
# 7. 메인 PPT 생성
# ──────────────────────────────────────────────
def generate_ppt(
    template_path: str,
    output_path: str,
    business_name: str,
    category: str,
    place_url: str,
    place_data: Dict[str, Any],
    keyword_stats: Dict[str, Any],
    main_keyword: str,
    rank: int,
):
    prs = Presentation(template_path)

    pc = keyword_stats.get("pc", 0)
    mobile = keyword_stats.get("mobile", 0)
    total = keyword_stats.get("total", 0)
    related = keyword_stats.get("related", [])
    rank_str = f"{rank}위" if rank > 0 else "순위권 외"

    # 의뢰 업체 정보 텍스트 구성
    address = place_data.get("address", "")
    phone = place_data.get("phone", "")
    biz_info_lines = [business_name]
    if category:
        biz_info_lines.append(category)
    if address:
        biz_info_lines.append(address)
    if phone:
        biz_info_lines.append(phone)
    biz_info = "\n".join(biz_info_lines)

    # 전체 공통 치환
    replacements = {
        "(업체명)": business_name,
        "한의원": category or "업종",
        "59위": rank_str,
        "Pc : 80회": f"PC : {pc:,}회",
        "모바일 : 80회": f"모바일 : {mobile:,}회",
        "전체 : 80회": f"전체 : {total:,}회",
        "플레이스 주소": biz_info,
        "상호명": business_name,
        "집중패키지": "집중패키지",
    }

    # 슬라이드2의 '집중패키지' (키워드 추천 칸)는 실제 메인 키워드로
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            text = " ".join(p.text for p in shape.text_frame.paragraphs).strip()
            if text == "집중패키지":
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("집중패키지", main_keyword)

    # 전체 치환 (집중패키지 제외한 나머지)
    replacements_no_package = {k: v for k, v in replacements.items() if k != "집중패키지"}
    replace_all_text(prs, replacements_no_package)

    # 슬라이드3 확장 키워드
    slide3 = prs.slides[2]
    if related:
        fill_slide3_keywords(slide3, related, business_name)

    # 슬라이드5 종합진단 업데이트
    slide5 = prs.slides[4]
    update_slide5_diagnosis(slide5, place_data)

    prs.save(output_path)
    print(f"[PPT] 저장 완료: {output_path}")


# ──────────────────────────────────────────────
# 8. 메인 실행
# ──────────────────────────────────────────────
async def main():
    print("=" * 60)
    print("네이버 플레이스 → PPT 자동 분석 시스템")
    print("=" * 60)

    # Step 1: 엑셀에서 첫 URL 행 읽기
    print("\n[Step 1] 엑셀 읽기...")
    row = read_first_url_row(EXCEL_PATH)
    if not row:
        print("❌ 네이버플레이스 URL이 있는 행을 찾지 못했습니다.")
        return

    business_name = str(row.get("업체명", "") or "")
    category = str(row.get("업종", "") or "")
    place_url = str(row.get("_place_url", "") or row.get("네이버플레이스", "") or row.get("주소", "") or "")
    place_id = extract_place_id(place_url)

    if not place_id:
        print(f"❌ URL에서 place_id를 추출하지 못했습니다: {place_url}")
        return

    print(f"  업체명: {business_name}")
    print(f"  업종: {category}")
    print(f"  place_id: {place_id}")

    # 메인 키워드 추론 (지역 + 업종)
    # URL에서 지역명 추출 시도
    city = ""
    decoded_url = urllib.parse.unquote(place_url)
    city_match = re.search(r"search/([가-힣]+)\s+[가-힣]+/place", decoded_url)
    if city_match:
        city = city_match.group(1)
    # 지역을 못 찾으면 URL의 searchText 파라미터에서
    if not city:
        st_match = re.search(r"searchText=([^&]+)", place_url)
        if st_match:
            st = urllib.parse.unquote(st_match.group(1))
            parts = st.split()
            if len(parts) >= 1:
                city = parts[0]

    main_keyword = f"{city} {category}".strip() if city and category else category or business_name
    print(f"  메인 키워드: {main_keyword}")

    # Step 2: Playwright 크롤링
    print(f"\n[Step 2] 네이버 플레이스 크롤링...")
    playwright = await async_playwright().start()
    browser = await playwright.chromium.launch(headless=True)

    try:
        place_data = await crawl_place(browser, place_id)

        # 업체명 업데이트 (크롤링으로 얻은 게 더 정확할 수 있음)
        if place_data.get("name"):
            business_name = place_data["name"]
        if place_data.get("category"):
            category = place_data["category"]

        # Step 3: 키워드 순위 조회
        print(f"\n[Step 3] '{main_keyword}' 순위 조회...")
        rank = await get_keyword_rank(browser, main_keyword, business_name)
        print(f"  현재 순위: {rank}위" if rank else "  순위: 미노출")

        # Step 4: 키워드 통계
        print(f"\n[Step 4] 키워드 통계 조회...")
        keyword_stats = await get_keyword_stats(main_keyword)

        # 연관 키워드가 부족하면 place_data의 키워드로 보완
        if len(keyword_stats.get("related", [])) < 4 and place_data.get("keywords"):
            extra = [{"keyword": k, "mobile": 0} for k in place_data["keywords"][:7]]
            keyword_stats["related"] = (keyword_stats.get("related", []) + extra)[:7]

    finally:
        await browser.close()
        await playwright.stop()

    # Step 5: PPT 생성
    print(f"\n[Step 5] PPT 생성...")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r'[\\/:*?"<>|]', "_", business_name)
    output_path = os.path.join(OUTPUT_DIR, f"{safe_name}_분석보고서_{timestamp}.pptx")

    generate_ppt(
        template_path=TEMPLATE_PATH,
        output_path=output_path,
        business_name=business_name,
        category=category,
        place_url=place_url,
        place_data=place_data,
        keyword_stats=keyword_stats,
        main_keyword=main_keyword,
        rank=rank,
    )

    print("\n" + "=" * 60)
    print("✅ 완료!")
    print(f"  업체: {business_name} ({category})")
    print(f"  순위: {rank}위" if rank else f"  순위: 미노출")
    print(f"  사진: {place_data['photo_count']}장")
    print(f"  방문자 리뷰: {place_data['visitor_review_count']}개")
    print(f"  블로그 리뷰: {place_data['blog_review_count']}개")
    print(f"  키워드 조회수: PC {keyword_stats['pc']:,} / 모바일 {keyword_stats['mobile']:,}")
    print(f"  PPT: {output_path}")
    print("=" * 60)


if __name__ == "__main__":
    if sys.platform == "win32":
        asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
    asyncio.run(main())
