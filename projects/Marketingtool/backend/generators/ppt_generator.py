# PPT Generator
"""
Generate PowerPoint proposal files using python-pptx
"""

from typing import List, Dict, Any, Optional
from datetime import datetime
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class PPTGenerator:
    """
    Generator for creating PowerPoint proposal files
    """

    def __init__(self):
        self.ppt = Presentation()
        self.title_font_size = Pt(32)
        self.heading_font_size = Pt(24)
        self.content_font_size = Pt(14)
        self.label_font_size = Pt(12)

        # Colors
        self.primary_color = RGBColor(221, 83, 53)    # #DD5335 (Primary)
        self.text_color = RGBColor(35, 35, 35)         # #232323 (Dark text)
        self.bg_color = RGBColor(250, 250, 250)      # #FAFAFA (Light bg)
        self.accent_color = RGBColor(59, 130, 246)    # #3B82F6 (Blue accent)
        self.success_color = RGBColor(34, 197, 94)      # Green
        self.warning_color = RGBColor(234, 179, 8)      # Yellow
        self.danger_color = RGBColor(239, 68, 68)       # Red

    def generate_proposal(
        self,
        analysis_data: Dict[str, Any],
        output_path: str = "./output"
    ) -> str:
        """
        Generate complete proposal PPT

        Args:
            analysis_data: All analysis data
            output_path: Output directory path

        Returns:
            Generated file path
        """
        request_id = analysis_data.get("id", "")
        business_name = analysis_data.get("business_name", "")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create output directory if not exists
        os.makedirs(output_path, exist_ok=True)

        # Generate slides
        self._add_cover_slide(analysis_data)  # Page 1
        self._add_exposure_slide(analysis_data)  # Page 2
        self._add_keyword_slide(analysis_data)   # Page 3
        self._add_diagnosis_slide(analysis_data) # Page 4

        # Save file
        filename = f"marketing_proposal_{request_id}_{timestamp}.pptx"
        filepath = os.path.join(output_path, filename)
        self.ppt.save(filepath)

        print(f"PPT generated: {filepath}")
        return filepath

    def _add_cover_slide(self, analysis_data: Dict[str, Any]):
        """Add Page 1: Business Overview"""
        slide = self.ppt.slides.add_slide()

        # Background
        background = slide.background
        background.fill.solid()
        background.fill.fg_color = self.primary_color

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].font.size = self.title_font_size
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.bold = True
        title_frame.text = "마케팅 분석 보고서"

        # Business name
        place_info = analysis_data.get("place_info", {})
        business_name = place_info.get("business_name", analysis_data.get("business_name", ""))

        name_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(9), Inches(0.5)
        )
        name_frame = name_box.text_frame
        name_frame.paragraphs[0].font.size = self.heading_font_size
        name_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        name_frame.paragraphs[0].font.bold = True
        name_frame.text = f"상호명: {business_name}"

        # Analysis summary
        keyword_stats = analysis_data.get("keyword_stats", {})
        current_rank = analysis_data.get("current_rank", {})

        y_pos = 3.5

        # Monthly search volume
        if keyword_stats:
            total_search = keyword_stats.get("monthly_search_pc", 0) + keyword_stats.get("monthly_search_mobile", 0)
            search_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos),
                Inches(9), Inches(0.4)
            )
            search_frame = search_box.text_frame
            search_frame.paragraphs[0].font.size = self.content_font_size
            search_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            search_frame.text = f"월간 검색량: {total_search:,}건 (PC: {keyword_stats.get('monthly_search_pc', 0):,} / Mobile: {keyword_stats.get('monthly_search_mobile', 0):,})"
            y_pos += 0.7

        # Current rank
        if current_rank:
            rank = current_rank.get("rank", -1)
            rank_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos),
                Inches(9), Inches(0.4)
            )
            rank_frame = rank_box.text_frame
            rank_frame.paragraphs[0].font.size = self.content_font_size
            rank_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            rank_text = f"현재 순위: {rank}위" if rank > 0 else "현재 순위: 노출 안됨"
            rank_frame.text = rank_text
            y_pos += 0.7

        # Date
        date_box = slide.shapes.add_textbox(
            Inches(7), Inches(6.5),
            Inches(2.5), Inches(0.3)
        )
        date_frame = date_box.text_frame
        date_frame.paragraphs[0].font.size = self.content_font_size
        date_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        date_frame.text = f"분석일: {datetime.now().strftime('%Y년 %m월 %d일')}"

    def _add_exposure_slide(self, analysis_data: Dict[str, Any]):
        """Add Page 2: Current Exposure Status"""
        slide = self.ppt.slides.add_slide()

        # Title
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.5)
        )
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = self.heading_font_size
        title_frame.paragraphs[0].font.bold = True
        title_frame.text = "현재 노출 상태"

        # Screenshot placeholder
        screenshot_path = analysis_data.get("current_rank_screenshot")
        if screenshot_path and os.path.exists(screenshot_path):
            # Add actual screenshot
            slide.shapes.add_picture(
                Inches(1), Inches(1),
                Inches(8), Inches(4),
                image_file=screenshot_path
            )
        else:
            # Placeholder box
            img_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(1),
                Inches(8), Inches(4)
            )
            fill = img_box.fill
            fill.solid()
            fill.fg_color = RGBColor(240, 240, 240)

            # Add note text
            note = slide.shapes.add_textbox(
                Inches(1), Inches(2.5),
                Inches(8), Inches(0.5)
            )
            note_frame = note.text_frame
            note_frame.paragraphs[0].font.size = Pt(18)
            note_frame.paragraphs[0].alignment = 1  # Center
            note_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            note_frame.text = "※ 스크린샷 준비 중"

        # Status summary boxes
        photo_info = analysis_data.get("photo_info", {})
        review_info = analysis_data.get("review_info", {})
        channel_info = analysis_data.get("channel_info", {})

        y_pos = 5.5

        # Photo status
        photo_text = f"대표 사진: {photo_info.get('photo_count', 0)}장"
        if photo_info.get("has_5_photos"):
            photo_text += " (좋음)"
            status_color = self.success_color
        else:
            photo_text += " (개선 필요)"
            status_color = self.warning_color

        photo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos),
            Inches(4), Inches(0.5)
        )
        photo_frame = photo_box.text_frame
        photo_frame.paragraphs[0].font.size = self.content_font_size
        photo_frame.paragraphs[0].font.color.rgb = status_color
        photo_frame.text = photo_text
        y_pos += 0.7

        # Review status
        review_text = f"사장님 답글: "
        review_text += "있음" if review_info.get("has_owner_reply") else "없음"
        review_box = slide.shapes.add_textbox(
            Inches(4.5), Inches(y_pos - 0.7),
            Inches(4.5), Inches(0.5)
        )
        review_frame = review_box.text_frame
        review_frame.paragraphs[0].font.size = self.content_font_size
        review_frame.paragraphs[0].font.color.rgb = (
            self.success_color if review_info.get("has_owner_reply") else self.danger_color
        )
        review_frame.text = review_text
        y_pos += 0.7

        # Channel status
        channel_text = f"외부 채널: "
        has_instagram = channel_info.get("has_instagram", False)
        has_kakao = channel_info.get("has_kakao_channel", False)

        if has_instagram or has_kakao:
            channel_text += "연동됨"
            status_color = self.success_color
        else:
            channel_text += "미연동"
            status_color = self.danger_color

        channel_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos),
            Inches(4), Inches(0.5)
        )
        channel_frame = channel_box.text_frame
        channel_frame.paragraphs[0].font.size = self.content_font_size
        channel_frame.paragraphs[0].font.color.rgb = status_color
        channel_frame.text = channel_text

    def _add_keyword_slide(self, analysis_data: Dict[str, Any]):
        """Add Page 3: Keyword Strategy"""
        slide = self.ppt.slides.add_slide()

        # Title
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.5)
        )
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = self.heading_font_size
        title_frame.paragraphs[0].font.bold = True
        title_frame.text = "키워드 확장 전략"

        # Keywords list
        keywords = analysis_data.get("expanded_keywords", [])
        y_pos = 1.2

        for i, kw in enumerate(keywords[:10]):  # Top 10 keywords
            keyword = kw.get("keyword", "")
            volume = kw.get("search_volume", 0)

            # Rank number
            rank_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos),
                Inches(0.5), Inches(0.35)
            )
            rank_frame = rank_box.text_frame
            rank_frame.paragraphs[0].font.size = Pt(16)
            rank_frame.paragraphs[0].font.bold = True
            rank_frame.paragraphs[0].font.color.rgb = self.primary_color
            rank_frame.text = f"{i + 1}"

            # Keyword
            kw_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y_pos),
                Inches(4), Inches(0.35)
            )
            kw_frame = kw_box.text_frame
            kw_frame.paragraphs[0].font.size = self.content_font_size
            kw_frame.text = keyword

            # Volume
            vol_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(y_pos),
                Inches(2.5), Inches(0.35)
            )
            vol_frame = vol_box.text_frame
            vol_frame.paragraphs[0].font.size = self.content_font_size
            vol_frame.paragraphs[0].alignment = 2  # Right align
            vol_text = f"{volume:,}건/월" if volume > 0 else "-"
            vol_frame.text = vol_text

            y_pos += 0.5

        if len(keywords) > 10:
            note = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos),
                Inches(9), Inches(0.3)
            )
            note_frame = note.text_frame
            note_frame.paragraphs[0].font.size = self.label_font_size
            note_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            note_frame.text = f"... 그 외 {len(keywords) - 10}개"

    def _add_diagnosis_slide(self, analysis_data: Dict[str, Any]):
        """Add Page 4: Diagnosis Results"""
        slide = self.ppt.slides.add_slide()

        # Title
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.5)
        )
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = self.heading_font_size
        title_frame.paragraphs[0].font.bold = True
        title_frame.text = "종합 진단 결과"

        # Diagnosis comments
        comments = analysis_data.get("diagnosis_comments", [])
        y_pos = 1.2

        for i, comment in enumerate(comments[:6]):  # Max 6 comments
            category = comment.get("category", "")
            message = comment.get("message", "")

            # Category badge
            category_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_pos),
                Inches(1.5), Inches(0.4)
            )
            cat_frame = category_box.text_frame
            cat_frame.paragraphs[0].font.size = self.label_font_size
            cat_frame.paragraphs[0].font.bold = True

            # Color by category
            if category == "review":
                cat_frame.paragraphs[0].font.color.rgb = self.danger_color
            elif category == "photo":
                cat_frame.paragraphs[0].font.color.rgb = self.warning_color
            else:
                cat_frame.paragraphs[0].font.color.rgb = self.accent_color

            cat_frame.text = f"[{category.upper()}]"

            # Comment
            msg_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(y_pos),
                Inches(7), Inches(0.4)
            )
            msg_frame = msg_box.text_frame
            msg_frame.paragraphs[0].font.size = self.content_font_size
            msg_frame.paragraphs[0].word_wrap = True
            msg_frame.text = message

            y_pos += 0.6

        # Footer note
        note_y = y_pos + 0.2
        note = slide.shapes.add_textbox(
            Inches(0.5), Inches(note_y),
            Inches(9), Inches(0.4)
        )
        note_frame = note.text_frame
        note_frame.paragraphs[0].font.size = self.label_font_size
        note_frame.paragraphs[0].font.color.rgb = self.primary_color
        note_frame.text = "위 항목을 개선하면 노출 상태가 향상될 수 있습니다."
